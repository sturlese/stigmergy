import datetime as dt
import hashlib
import io
import shutil
import struct
import zlib

import pymupdf
import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches

from stigmergy.capture import artifacts, schema
from stigmergy.capture import extraction as extraction_module
from stigmergy.capture.errors import ArtifactRejected, ExtractionError
from stigmergy.capture.evidence import MemoryEvidenceStore
from stigmergy.capture.extraction import ExtractionResult, extract_artifact


def _digital_pdf(text: str) -> bytes:
    document = pymupdf.open()
    page = document.new_page()
    page.insert_text((72, 100), text, fontsize=14)
    data = document.tobytes()
    document.close()
    return data


def _image_bytes(text: str, image_type: str = "png") -> bytes:
    document = pymupdf.open(stream=_digital_pdf(text), filetype="pdf")
    pixmap = document[0].get_pixmap(dpi=200, alpha=False)
    data = pixmap.tobytes(image_type)
    document.close()
    return data


def _scanned_pdf(text: str) -> bytes:
    image = _image_bytes(text)
    document = pymupdf.open()
    page = document.new_page()
    page.insert_image(page.rect, stream=image)
    data = document.tobytes()
    document.close()
    return data


def test_plain_text_extraction_is_byte_exact():
    data = b"First line\nSecond line\n"
    result = extract_artifact(data, schema.MEDIA_TEXT)

    assert result.text.encode() == data
    assert result.ocr_pages == ()


def test_readable_output_limit_is_enforced_before_child_serialization(monkeypatch):
    monkeypatch.setattr(extraction_module, "MAX_CAPTURE_EXTRACTED_BYTES", 10)

    with pytest.raises(ExtractionError, match="readable byte limit"):
        extract_artifact(b"more than ten bytes", schema.MEDIA_TEXT)


def test_capture_extraction_rejects_aggregate_readable_bytes_without_derivatives(monkeypatch):
    store = MemoryEvidenceStore()

    def artifact(data):
        digest = hashlib.sha256(data).hexdigest()
        return schema.ArtifactRef(
            blob_ref=store.put(data),
            sha256=digest,
            bytes=len(data),
            media_type=schema.MEDIA_PDF,
        )

    first = artifact(b"first")
    second = artifact(b"second")
    envelope = schema.CaptureEnvelope(
        idempotency_key="aggregate-readable",
        actor=schema.Actor(subject="alice", display_name="Alice"),
        audience=None,
        origin=schema.Origin(adapter="mcp", captured_at=dt.datetime.now(dt.UTC)),
        artifacts=(first, second),
    )
    original_refs = set(store.objects)
    monkeypatch.setattr(extraction_module, "MAX_CAPTURE_EXTRACTED_BYTES", 10, raising=False)
    monkeypatch.setattr(
        extraction_module,
        "extract_artifact",
        lambda *_args, **_kwargs: ExtractionResult(
            text="abcdef",
            media_type=schema.MEDIA_PDF,
            extractor="fixture",
        ),
    )

    with pytest.raises(ExtractionError):
        extraction_module.extract_capture(store, envelope, bounded=False)

    assert set(store.objects) == original_refs


def test_html_extraction_keeps_title_and_readable_content():
    data = (
        b"<html><head><title>Plan</title><style>hidden</style></head>"
        b"<body><nav>menu</nav><main><h1>Decision</h1><p>Ship it.</p></main></body></html>"
    )
    result = extract_artifact(data, schema.MEDIA_HTML)

    assert result.text == "# Plan\n\nDecision\n\nShip it."
    assert "hidden" not in result.text
    assert "menu" not in result.text


def test_digital_pdf_does_not_use_ocr():
    result = extract_artifact(
        _digital_pdf("Digital text is complete and long enough for extraction."),
        schema.MEDIA_PDF,
    )

    assert "Digital text is complete" in result.text
    assert result.ocr_pages == ()
    assert result.decisions == ("page 1: digital text",)


@pytest.mark.skipif(shutil.which("tesseract") is None, reason="tesseract is not installed")
def test_scanned_pdf_uses_ocr_and_records_the_page():
    result = extract_artifact(
        _scanned_pdf("Scanned quarterly roadmap"),
        schema.MEDIA_PDF,
    )

    assert "Scanned quarterly roadmap" in result.text
    assert result.ocr_pages == (1,)
    assert result.decisions == ("page 1: OCR",)


def test_docx_preserves_heading_paragraph_and_table_order():
    document = Document()
    document.add_heading("Roadmap", level=1)
    document.add_paragraph("First milestone")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Owner"
    table.cell(0, 1).text = "Alice"
    output = io.BytesIO()
    document.save(output)

    result = extract_artifact(output.getvalue(), schema.MEDIA_DOCX)

    assert result.text == "# Roadmap\n\nFirst milestone\n\n| Owner | Alice |"


def test_pptx_preserves_slide_title_body_table_and_notes():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Launch"
    slide.placeholders[1].text = "Approve the rollout"
    table = slide.shapes.add_table(1, 2, Inches(1), Inches(3), Inches(5), Inches(1)).table
    table.cell(0, 0).text = "Owner"
    table.cell(0, 1).text = "Bob"
    slide.notes_slide.notes_text_frame.text = "Mention the Friday deadline"
    output = io.BytesIO()
    presentation.save(output)

    result = extract_artifact(output.getvalue(), schema.MEDIA_PPTX)

    assert "## Slide 1" in result.text
    assert "### Launch" in result.text
    assert "Approve the rollout" in result.text
    assert "| Owner | Bob |" in result.text
    assert "### Speaker notes" in result.text
    assert "Mention the Friday deadline" in result.text


@pytest.mark.skipif(shutil.which("tesseract") is None, reason="tesseract is not installed")
@pytest.mark.parametrize(
    ("image_type", "media_type"),
    [("png", schema.MEDIA_PNG), ("jpeg", schema.MEDIA_JPEG)],
)
def test_image_formats_use_ocr(image_type, media_type):
    result = extract_artifact(_image_bytes("Image decision text", image_type), media_type)

    assert "Image decision text" in result.text
    assert result.ocr_pages == (1,)


def test_image_dimensions_are_bounded_before_raster_decode():
    data = bytearray(_image_bytes("small image"))
    data[16:20] = struct.pack(">I", 10_000)
    data[20:24] = struct.pack(">I", 10_000)
    data[29:33] = struct.pack(">I", zlib.crc32(data[12:29]) & 0xFFFFFFFF)

    with pytest.raises(ArtifactRejected, match="pixel safety limit"):
        extract_artifact(bytes(data), schema.MEDIA_PNG)


def test_pdf_ocr_dimensions_are_bounded_before_page_rasterization():
    document = pymupdf.open()
    document.new_page(width=100_000, height=100_000)
    data = document.tobytes()
    document.close()

    with pytest.raises(ArtifactRejected, match="OCR pixel safety limit"):
        extract_artifact(data, schema.MEDIA_PDF)


def test_mime_spoofing_fails_before_parser_execution():
    with pytest.raises(ArtifactRejected, match="does not match"):
        extract_artifact(b"not a PDF", schema.MEDIA_PDF)


def test_openxml_detection_identifies_real_packages():
    document = Document()
    document.add_paragraph("Text")
    output = io.BytesIO()
    document.save(output)

    assert artifacts.detect_media(output.getvalue()) == schema.MEDIA_DOCX
