import asyncio
import datetime as dt
import io
import json

import httpx
from docx import Document
from pptx import Presentation

from stigmergy.bridge import drive as drive_module
from stigmergy.bridge.acquire import Acquirer
from stigmergy.bridge.cloud import AcquiredArtifact, CloudClient
from stigmergy.bridge.drive import (
    GOOGLE_DOC,
    GOOGLE_SLIDES,
    SCOPES,
    DriveClient,
    file_id_from_url,
)
from stigmergy.bridge.server import build_mcp
from stigmergy.capture import schema
from stigmergy.capture.extraction import extract_artifact


class FakeKeyring:
    def __init__(self, value=None):
        self.value = value
        self.saved = []

    def get_password(self, service, account):
        return self.value

    def set_password(self, service, account, value):
        self.saved.append((service, account, value))


def _authorized_credentials():
    return json.dumps(
        {
            "token": "local-access-token",
            "refresh_token": "local-refresh-token",
            "token_uri": "https://oauth2.googleapis.com/token",
            "client_id": "local-client-id",
            "client_secret": "local-client-secret",
            "scopes": list(SCOPES),
            "expiry": "2099-01-01T00:00:00Z",
        }
    )


def test_drive_credentials_round_trip_through_keychain_only():
    keychain = FakeKeyring(_authorized_credentials())
    credentials = DriveClient("", keyring_module=keychain)._credentials()
    assert credentials.refresh_token == "local-refresh-token"
    assert len(keychain.saved) == 1
    assert "local-refresh-token" in keychain.saved[0][2]


def test_drive_url_parser_accepts_files_docs_and_slides():
    assert file_id_from_url("https://drive.google.com/file/d/file_1/view") == "file_1"
    assert file_id_from_url("https://docs.google.com/document/d/doc-2/edit") == "doc-2"
    assert file_id_from_url("https://docs.google.com/presentation/d/slides_3/edit") == "slides_3"
    assert file_id_from_url("https://example.com/file/d/no") is None


class _DriveRequest:
    def __init__(self, result=None):
        self.result = result

    def execute(self):
        return self.result


class _DriveFiles:
    def __init__(self, metadata):
        self.metadata = metadata
        self.export_type = None

    def get(self, **_kwargs):
        return _DriveRequest(self.metadata)

    def export_media(self, *, mimeType, **_kwargs):
        self.export_type = mimeType
        return _DriveRequest()

    def get_media(self, **_kwargs):
        return _DriveRequest()


class _DriveService:
    def __init__(self, metadata):
        self.resource = _DriveFiles(metadata)

    def files(self):
        return self.resource


def _docx_export() -> bytes:
    document = Document()
    document.add_heading("Exported plan", level=1)
    document.add_paragraph("First decision")
    document.add_paragraph("Second decision")
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _pptx_export() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Exported launch"
    slide.placeholders[1].text = "First milestone\nSecond milestone"
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def test_google_doc_export_flows_into_ordered_docx_extraction(monkeypatch):
    data = _docx_export()
    service = _DriveService({"id": "doc-2", "name": "Plan", "mimeType": GOOGLE_DOC})
    monkeypatch.setattr(drive_module, "_download", lambda _request: data)

    artifact = DriveClient("", service=service).acquire(
        "https://docs.google.com/document/d/doc-2/edit"
    )
    extracted = extract_artifact(artifact.data, artifact.media_type)

    assert service.resource.export_type == schema.MEDIA_DOCX
    assert artifact.original_name == "Plan.docx"
    assert artifact.acquisition is not None
    assert artifact.acquisition.export_media_type == schema.MEDIA_DOCX
    assert extracted.text.index("Exported plan") < extracted.text.index("First decision")
    assert extracted.text.index("First decision") < extracted.text.index("Second decision")


def test_drive_provenance_is_canonical_and_drops_resource_keys(monkeypatch):
    data = _docx_export()
    service = _DriveService({"id": "doc-2", "name": "Plan", "mimeType": GOOGLE_DOC})
    monkeypatch.setattr(drive_module, "_download", lambda _request: data)

    artifact = DriveClient("", service=service).acquire(
        "https://docs.google.com/document/d/doc-2/edit?resourcekey=secret"
    )

    assert artifact.source_url == "https://drive.google.com/file/d/doc-2/view"
    assert artifact.locator == artifact.source_url
    assert "secret" not in json.dumps(artifact.acquisition.model_dump(mode="json"))
    assert artifact.acquisition.original_url == (
        "https://docs.google.com/document/d/doc-2/edit"
    )


def test_google_slides_export_flows_into_ordered_pptx_extraction(monkeypatch):
    data = _pptx_export()
    service = _DriveService(
        {"id": "slides-3", "name": "Launch", "mimeType": GOOGLE_SLIDES}
    )
    monkeypatch.setattr(drive_module, "_download", lambda _request: data)

    artifact = DriveClient("", service=service).acquire(
        "https://docs.google.com/presentation/d/slides-3/edit"
    )
    extracted = extract_artifact(artifact.data, artifact.media_type)

    assert service.resource.export_type == schema.MEDIA_PPTX
    assert artifact.original_name == "Launch.pptx"
    assert artifact.acquisition is not None
    assert artifact.acquisition.export_media_type == schema.MEDIA_PPTX
    assert extracted.text.index("Exported launch") < extracted.text.index("First milestone")
    assert extracted.text.index("First milestone") < extracted.text.index("Second milestone")


def test_cloud_upload_uses_raw_bytes_and_never_sends_google_credentials():
    seen = []
    uploaded = b"%PDF-1.7\n" + b"x" * (2 * 1024 * 1024)

    def handler(request):
        seen.append(request)
        if request.url.host == "upload.example":
            assert request.content == uploaded
            assert request.headers["content-length"] == str(len(uploaded))
            return httpx.Response(200)
        body = json.loads(request.content)
        serialized = request.content.decode()
        assert "local-refresh-token" not in serialized
        assert "local-access-token" not in serialized
        if request.url.path == "/bridge/uploads":
            assert "base64" not in serialized
            return httpx.Response(
                200,
                json={
                    "upload_id": "e19670ca-57b1-4e64-a412-a2027f461266",
                    "upload_url": "https://upload.example/object",
                    "expires_at": "2026-08-24T12:00:00+00:00",
                },
            )
        assert request.url.path == "/bridge/captures"
        assert body["acquisition"]["drive_file_id"] == "file_1"
        return httpx.Response(
            200,
            json={"id": "capture-1", "status": "queued", "created": True},
        )

    client = httpx.Client(transport=httpx.MockTransport(handler))
    cloud = CloudClient("https://brain.example", "member-token", client=client)
    receipt = cloud.submit_artifacts(
        [
            AcquiredArtifact(
                data=uploaded,
                media_type=schema.MEDIA_PDF,
                original_name="board.pdf",
                source_url="https://drive.google.com/file/d/file_1/view",
                locator="https://drive.google.com/file/d/file_1/view",
                acquisition=schema.AcquisitionProvenance(
                    original_url="https://drive.google.com/file/d/file_1/view",
                    final_url="https://drive.google.com/file/d/file_1/view",
                    drive_file_id="file_1",
                    drive_media_type=schema.MEDIA_PDF,
                    acquired_at=dt.datetime(2026, 8, 24, tzinfo=dt.UTC),
                ),
            )
        ],
        title="Board pack",
        occurred_at=None,
        audience=["finance"],
    )
    assert receipt == {"id": "capture-1", "status": "queued", "created": True}
    assert [request.method for request in seen] == ["POST", "PUT", "POST"]
    assert [len(request.content) for request in seen if request.method == "PUT"] == [len(uploaded)]
    assert all(len(request.content) < 16_384 for request in seen if request.method == "POST")


def test_private_drive_submission_keeps_google_credentials_inside_the_local_bridge(monkeypatch):
    keychain = FakeKeyring(_authorized_credentials())
    exported = _docx_export()
    service = _DriveService({"id": "doc-2", "name": "Plan", "mimeType": GOOGLE_DOC})
    local_credentials = []

    def local_build(_api, _version, *, credentials, cache_discovery):
        local_credentials.append(credentials)
        assert cache_discovery is False
        return service

    monkeypatch.setattr(drive_module, "build", local_build)
    monkeypatch.setattr(drive_module, "_download", lambda _request: exported)
    requests = []

    def handler(request):
        requests.append(request)
        wire = request.content + b"\n" + b"\n".join(
            f"{key}: {value}".encode() for key, value in request.headers.items()
        )
        assert b"local-refresh-token" not in wire
        assert b"local-access-token" not in wire
        if request.url.host == "upload.example":
            assert request.content == exported
            return httpx.Response(200)
        body = json.loads(request.content)
        if request.url.path == "/bridge/uploads":
            return httpx.Response(
                200,
                json={
                    "upload_id": "e19670ca-57b1-4e64-a412-a2027f461266",
                    "upload_url": "https://upload.example/object",
                    "expires_at": "2026-08-24T12:00:00+00:00",
                },
            )
        assert body["acquisition"]["drive_file_id"] == "doc-2"
        return httpx.Response(
            200,
            json={"id": "capture-1", "status": "queued", "created": True},
        )

    cloud = CloudClient(
        "https://brain.example",
        "member-token",
        client=httpx.Client(transport=httpx.MockTransport(handler)),
    )
    mcp = build_mcp(
        cloud,
        Acquirer(DriveClient("", keyring_module=keychain)),
    )
    blocks, _ = asyncio.run(
        mcp.call_tool(
            "brain_submit",
            {
                "url": "https://docs.google.com/document/d/doc-2/edit?resourcekey=private",
                "audience": ["finance"],
            },
        )
    )

    assert json.loads(blocks[0].text)["status"] == "queued"
    assert len(local_credentials) == 1
    assert local_credentials[0].refresh_token == "local-refresh-token"
    assert [request.method for request in requests] == ["POST", "PUT", "POST"]


def test_acquirer_keeps_text_and_local_file_exact(tmp_path):
    text = "Decision: ship café.\n"
    assert Acquirer().text(text).data == text.encode("utf-8")
    path = tmp_path / "note.md"
    path.write_bytes(text.encode("utf-8"))
    artifact = Acquirer().path(str(path))
    assert artifact.data == text.encode("utf-8")
    assert artifact.media_type == schema.MEDIA_MARKDOWN


def test_local_tool_rejects_zero_or_multiple_inputs():
    class NoCloud:
        async def call_tool(self, name, arguments):
            raise AssertionError("cloud should not be called")

    mcp = build_mcp(NoCloud(), Acquirer())

    async def call(arguments):
        blocks, _ = await mcp.call_tool("brain_submit", arguments)
        return json.loads(blocks[0].text)

    assert asyncio.run(call({})) == {"error": "provide exactly one of text, path, or url"}
    assert asyncio.run(call({"text": "a", "path": "/tmp/a"})) == {
        "error": "provide exactly one of text, path, or url"
    }


def test_local_search_forwards_only_the_public_search_arguments():
    cloud = CloudClient("https://brain.example", "member-token")
    forwarded = {}

    async def call_tool(name, arguments):
        forwarded["name"] = name
        forwarded["arguments"] = arguments
        return json.dumps({"hits": []})

    cloud.call_tool = call_tool
    mcp = build_mcp(cloud, Acquirer())

    blocks, _ = asyncio.run(
        mcp.call_tool(
            "search_brain",
            {
                "query": "renewal cadence",
                "filters": {"entity": "acme"},
                "max_results": 3,
            },
        )
    )

    assert json.loads(blocks[0].text) == {"hits": []}
    assert forwarded == {
        "name": "search_brain",
        "arguments": {
            "query": "renewal cadence",
            "filters": {"entity": "acme"},
            "max_results": 3,
        },
    }


def test_local_bridge_forwards_every_non_acquisition_tool_without_changing_arguments():
    class RecordingCloud:
        def __init__(self):
            self.calls = []

        async def call_tool(self, name, arguments):
            self.calls.append((name, arguments))
            return json.dumps({"tool": name})

    cloud = RecordingCloud()
    mcp = build_mcp(cloud, Acquirer())
    cases = [
        ("read_page", {"path": "wiki/notes/Terms.md"}),
        ("list_entities", {}),
        ("describe_entity", {"entity": "Acme"}),
        ("ask", {"question": "What is the renewal term?"}),
        ("brain_submissions", {"limit": 7, "status": "failed"}),
        (
            "brain_delete",
            {"paths": ["wiki/notes/Obsolete.md"], "why": "Superseded"},
        ),
    ]

    async def call_all():
        responses = []
        for name, arguments in cases:
            blocks, _ = await mcp.call_tool(name, arguments)
            responses.append(json.loads(blocks[0].text))
        return responses

    responses = asyncio.run(call_all())

    assert cloud.calls == cases
    assert responses == [{"tool": name} for name, _arguments in cases]


def test_local_bridge_acquires_text_and_paths_before_using_the_shared_upload_flow(tmp_path):
    class UploadCloud:
        def __init__(self):
            self.uploads = []

        def submit_artifacts(self, artifacts, **metadata):
            self.uploads.append((artifacts, metadata))
            return {"id": f"capture-{len(self.uploads)}", "status": "queued"}

    path = tmp_path / "decision.txt"
    path.write_bytes(b"Path decision")
    cloud = UploadCloud()
    mcp = build_mcp(cloud, Acquirer())

    async def submit(arguments):
        blocks, _ = await mcp.call_tool("brain_submit", arguments)
        return json.loads(blocks[0].text)

    assert asyncio.run(submit({"text": "Text decision"}))["id"] == "capture-1"
    assert asyncio.run(submit({"path": str(path)}))["id"] == "capture-2"
    assert cloud.uploads[0][0][0].data == b"Text decision"
    assert cloud.uploads[1][0][0].data == b"Path decision"
