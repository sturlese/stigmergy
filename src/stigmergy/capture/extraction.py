"""Deterministic readable extraction with bounded parser and OCR execution."""

from __future__ import annotations

import datetime as dt
import io
import math
import multiprocessing
import re
import sys
import time
import warnings
from collections.abc import Callable
from dataclasses import dataclass
from html.parser import HTMLParser

from pydantic import BaseModel, ConfigDict

from stigmergy.capture import artifacts, schema
from stigmergy.capture import evidence as evidence_module
from stigmergy.capture.errors import ArtifactRejected, CaptureError, ExtractionError
from stigmergy.kernel.deadline import hard_deadline
from stigmergy.kernel.llm import OCR_MODEL

EXTRACTOR_VERSION = "1"
MAX_PAGES = 200
MAX_IMAGE_PIXELS = 50_000_000
MAX_CAPTURE_EXTRACTED_BYTES = 2 * 1024 * 1024
DEFAULT_TIMEOUT_S = 120
CAPTURE_TIMEOUT_S = 180
OCR_DPI = 200
MAX_CHILD_ADDRESS_SPACE_BYTES = 1024 * 1024 * 1024


class ExtractionResult(BaseModel):
    model_config = ConfigDict(frozen=True, extra="forbid")

    text: str
    media_type: str
    extractor: str
    extractor_version: str = EXTRACTOR_VERSION
    pages: int = 1
    ocr_pages: tuple[int, ...] = ()
    decisions: tuple[str, ...] = ()


@dataclass(frozen=True)
class ExtractedArtifact:
    original: schema.ArtifactRef
    readable_ref: str
    readable_sha256: str
    readable_bytes: int
    result: ExtractionResult

    def metadata(self) -> dict:
        return {
            "original": self.original.model_dump(mode="json"),
            "readable_ref": self.readable_ref,
            "readable_sha256": self.readable_sha256,
            "readable_bytes": self.readable_bytes,
            "extractor": self.result.extractor,
            "extractor_version": self.result.extractor_version,
            "pages": self.result.pages,
            "ocr_pages": list(self.result.ocr_pages),
            "decisions": list(self.result.decisions),
        }


def extract_artifact(
    data: bytes,
    media_type: str,
    *,
    ocr_model: str = OCR_MODEL,
    vision_ocr: Callable[[bytes, str], str] | None = None,
) -> ExtractionResult:
    _validate_ocr_model(ocr_model)
    detected = artifacts.detect_media(data, declared=media_type)
    if detected in {schema.MEDIA_TEXT, schema.MEDIA_MARKDOWN}:
        result = _text(data, detected)
    elif detected == schema.MEDIA_HTML:
        result = _html(data)
    elif detected == schema.MEDIA_PDF:
        result = _pdf(
            data,
            ocr_model=ocr_model,
            vision_ocr=vision_ocr,
        )
    elif detected == schema.MEDIA_DOCX:
        result = _docx(data)
    elif detected == schema.MEDIA_PPTX:
        result = _pptx(data)
    elif detected in {schema.MEDIA_PNG, schema.MEDIA_JPEG}:
        result = _image(
            data,
            detected,
            ocr_model=ocr_model,
            vision_ocr=vision_ocr,
        )
    elif detected == schema.MEDIA_SLACK:
        result = _slack(data)
    else:
        raise ArtifactRejected("unsupported artifact format")
    _guard_readable(result.text)
    return result


def extract_bounded(
    data: bytes,
    media_type: str,
    *,
    timeout_s: int = DEFAULT_TIMEOUT_S,
    ocr_model: str = OCR_MODEL,
) -> ExtractionResult:
    context = multiprocessing.get_context("spawn")
    parent, child = context.Pipe(duplex=False)
    process = context.Process(
        target=_extract_child,
        args=(
            child,
            data,
            media_type,
            ocr_model,
            max(1, int(timeout_s)),
        ),
        daemon=True,
    )
    process.start()
    child.close()
    try:
        if not parent.poll(max(1, int(timeout_s))):
            process.terminate()
            process.join(timeout=5)
            raise ExtractionError("artifact extraction timed out")
        try:
            outcome = parent.recv()
        except (EOFError, OSError) as error:
            raise ExtractionError("artifact extraction process failed") from error
    finally:
        parent.close()
        if process.is_alive():
            process.terminate()
        process.join(timeout=5)
    if outcome["ok"]:
        return ExtractionResult.model_validate(outcome["result"])
    if outcome["category"] == ArtifactRejected.category:
        raise ArtifactRejected(outcome["error"])
    raise ExtractionError(outcome["error"])


def _extract_child(
    connection,
    data: bytes,
    media_type: str,
    ocr_model: str,
    timeout_s: int,
) -> None:
    try:
        _apply_resource_limits(timeout_s)
        result = extract_artifact(
            data,
            media_type,
            ocr_model=ocr_model,
        )
        connection.send({"ok": True, "result": result.model_dump(mode="json")})
    except CaptureError as error:
        connection.send(
            {"ok": False, "category": error.category, "error": str(error)}
        )
    except Exception as error:
        connection.send(
            {
                "ok": False,
                "category": ExtractionError.category,
                "error": f"artifact extraction failed ({error.__class__.__name__})",
            }
        )
    finally:
        connection.close()


def extract_capture(
    store,
    envelope: schema.CaptureEnvelope,
    *,
    bounded: bool = True,
    timeout_s: int = CAPTURE_TIMEOUT_S,
    ocr_model: str = OCR_MODEL,
    vision_ocr: Callable[[bytes, str], str] | None = None,
) -> tuple[ExtractedArtifact, ...]:
    _validate_ocr_model(ocr_model)
    if bounded and vision_ocr is not None:
        raise ExtractionError("custom OCR handlers require unbounded extraction")
    with hard_deadline(
        timeout_s if bounded else None,
        lambda: ExtractionError("capture extraction timed out"),
    ):
        return _extract_capture(
            store,
            envelope,
            bounded=bounded,
            timeout_s=timeout_s,
            ocr_model=ocr_model,
            vision_ocr=vision_ocr,
        )


def _extract_capture(
    store,
    envelope: schema.CaptureEnvelope,
    *,
    bounded: bool,
    timeout_s: int,
    ocr_model: str,
    vision_ocr: Callable[[bytes, str], str] | None,
) -> tuple[ExtractedArtifact, ...]:
    pending = []
    readable_total = 0
    deadline = time.monotonic() + max(1, int(timeout_s)) if bounded else None

    def remaining() -> int:
        if deadline is None:
            return DEFAULT_TIMEOUT_S
        seconds = deadline - time.monotonic()
        if seconds <= 0:
            raise ExtractionError("capture extraction timed out")
        return min(DEFAULT_TIMEOUT_S, max(1, math.ceil(seconds)))

    for artifact in envelope.artifacts:
        remaining()
        if store.head(artifact.blob_ref).bytes != artifact.bytes:
            raise ExtractionError("original artifact failed its digest or size check")
        data = store.get_limited(artifact.blob_ref, max_bytes=artifact.bytes)
        remaining()
        if (
            len(data) != artifact.bytes
            or evidence_module.sha256(data) != artifact.sha256
        ):
            raise ExtractionError("original artifact failed its digest or size check")
        extract = extract_bounded if bounded else extract_artifact
        result = extract(
            data,
            artifact.media_type,
            **(
                {
                    "timeout_s": remaining(),
                    "ocr_model": ocr_model,
                }
                if bounded
                else {"ocr_model": ocr_model, "vision_ocr": vision_ocr}
            ),
        )
        remaining()
        readable = result.text.encode("utf-8")
        if not readable:
            raise ExtractionError("readable extraction is empty")
        readable_total += len(readable)
        if readable_total > MAX_CAPTURE_EXTRACTED_BYTES:
            raise ExtractionError("extraction exceeds the capture-wide readable byte limit")
        pending.append((artifact, readable, result))

    extracted = []
    for artifact, readable, result in pending:
        remaining()
        readable_ref = (
            artifact.blob_ref
            if artifact.media_type in {schema.MEDIA_TEXT, schema.MEDIA_MARKDOWN}
            else store.put(readable)
        )
        remaining()
        extracted.append(
            ExtractedArtifact(
                original=artifact,
                readable_ref=readable_ref,
                readable_sha256=evidence_module.sha256(readable),
                readable_bytes=len(readable),
                result=result,
            )
        )
    return tuple(extracted)


def _text(data: bytes, media_type: str) -> ExtractionResult:
    if len(data) > MAX_CAPTURE_EXTRACTED_BYTES:
        raise ExtractionError("extraction exceeds the readable byte limit")
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError as error:
        raise ArtifactRejected("text artifact is not valid UTF-8") from error
    return ExtractionResult(text=text, media_type=media_type, extractor="utf8")


class _ReadableHTML(HTMLParser):
    _BLOCKS = {
        "article",
        "blockquote",
        "br",
        "dd",
        "div",
        "dl",
        "dt",
        "figcaption",
        "h1",
        "h2",
        "h3",
        "h4",
        "h5",
        "h6",
        "li",
        "main",
        "p",
        "pre",
        "section",
        "table",
        "td",
        "th",
        "tr",
    }
    _IGNORED = {"script", "style", "svg", "canvas", "noscript", "nav", "footer"}

    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.depth = 0
        self.title_depth = 0
        self.title: list[str] = []
        self.parts: list[str] = []
        self.readable_bytes = 0

    def append(self, target: list[str], value: str) -> None:
        self.readable_bytes += len(value.encode("utf-8"))
        if self.readable_bytes > MAX_CAPTURE_EXTRACTED_BYTES:
            raise ExtractionError("extraction exceeds the readable byte limit")
        target.append(value)

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag in self._IGNORED:
            self.depth += 1
        if tag == "title":
            self.title_depth += 1
        if not self.depth and tag in self._BLOCKS:
            self.append(self.parts, "\n")

    def handle_endtag(self, tag: str) -> None:
        if tag == "title" and self.title_depth:
            self.title_depth -= 1
        if tag in self._IGNORED and self.depth:
            self.depth -= 1
        if not self.depth and tag in self._BLOCKS:
            self.append(self.parts, "\n")

    def handle_data(self, data: str) -> None:
        if self.title_depth:
            self.append(self.title, data)
        elif not self.depth:
            self.append(self.parts, data)


def _html(data: bytes) -> ExtractionResult:
    parser = _ReadableHTML()
    try:
        parser.feed(data.decode("utf-8"))
        parser.close()
    except (UnicodeDecodeError, ValueError) as error:
        raise ArtifactRejected("HTML artifact is not valid UTF-8") from error
    body = _clean_blocks("".join(parser.parts))
    title = " ".join("".join(parser.title).split())
    text = f"# {title}\n\n{body}" if title and body else title or body
    if not text:
        raise ExtractionError("HTML artifact contains no readable text")
    return ExtractionResult(text=text, media_type=schema.MEDIA_HTML, extractor="html")


def _pdf(
    data: bytes,
    *,
    ocr_model: str,
    vision_ocr: Callable[[bytes, str], str] | None,
) -> ExtractionResult:
    import pymupdf

    try:
        document = pymupdf.open(stream=data, filetype="pdf")
    except Exception as error:
        raise ArtifactRejected("PDF is corrupt") from error
    with document:
        if document.needs_pass:
            raise ArtifactRejected("encrypted PDFs are not supported")
        if document.page_count > MAX_PAGES:
            raise ArtifactRejected(f"PDF exceeds the {MAX_PAGES}-page limit")
        pages = _ReadableBuilder()
        ocr_pages = []
        decisions = []
        for index, page in enumerate(document, start=1):
            text = page.get_text("text", sort=True)
            if _needs_ocr(text):
                _guard_pdf_ocr_geometry(page.rect.width, page.rect.height)
                try:
                    handler = _vision_handler(ocr_model, vision_ocr)
                    pixmap = page.get_pixmap(dpi=OCR_DPI, alpha=False)
                    ocr_text = handler(pixmap.tobytes("png"), schema.MEDIA_PNG)
                except Exception as error:
                    raise ExtractionError(
                        f"OCR failed for PDF page {index} ({error.__class__.__name__})"
                    ) from error
                if ocr_text.strip():
                    text = ocr_text
                ocr_pages.append(index)
                decisions.append(f"page {index}: OCR")
            else:
                decisions.append(f"page {index}: digital text")
            pages.add(f"## Page {index}\n\n{text.rstrip()}")
        combined = pages.text().rstrip()
        if not combined:
            raise ExtractionError("PDF contains no readable text")
        return ExtractionResult(
            text=combined,
            media_type=schema.MEDIA_PDF,
            extractor=(
                f"pymupdf+{ocr_model.removeprefix('openrouter:')}"
                if ocr_pages
                else "pymupdf"
            ),
            pages=document.page_count,
            ocr_pages=tuple(ocr_pages),
            decisions=tuple(decisions),
        )


def _needs_ocr(text: str) -> bool:
    stripped = text.strip()
    if sum(char.isalnum() for char in stripped) < 20:
        return True
    replacements = stripped.count("\ufffd")
    return replacements > max(2, len(stripped) // 50)


def _docx(data: bytes) -> ExtractionResult:
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    artifacts.validate_openxml(data)
    try:
        document = Document(io.BytesIO(data))
    except Exception as error:
        raise ArtifactRejected("DOCX is corrupt") from error
    blocks = _ReadableBuilder()
    for block in document.iter_inner_content():
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if not text:
                continue
            style = str(block.style.name or "")
            match = re.match(r"Heading\s+([1-6])$", style, re.I)
            blocks.add(f"{'#' * int(match.group(1))} {text}" if match else text)
        elif isinstance(block, Table):
            for row in _table_rows(block.rows):
                blocks.add(row)
    text = blocks.text().strip()
    if not text:
        raise ExtractionError("DOCX contains no readable text")
    return ExtractionResult(text=text, media_type=schema.MEDIA_DOCX, extractor="python-docx")


def _table_rows(rows) -> list[str]:
    result = []
    for row in rows:
        cells = [" ".join(cell.text.split()) for cell in row.cells]
        result.append("| " + " | ".join(cells) + " |")
    return result


def _pptx(data: bytes) -> ExtractionResult:
    from pptx import Presentation

    artifacts.validate_openxml(data)
    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as error:
        raise ArtifactRejected("PPTX is corrupt") from error
    if len(presentation.slides) > MAX_PAGES:
        raise ArtifactRejected(f"PPTX exceeds the {MAX_PAGES}-slide limit")
    slides = _ReadableBuilder()
    for index, slide in enumerate(presentation.slides, start=1):
        lines = [f"## Slide {index}"]
        title_shape = slide.shapes.title
        if title_shape is not None and title_shape.text.strip():
            lines.append(f"### {title_shape.text.strip()}")
        shapes = sorted(
            (shape for shape in slide.shapes if shape is not title_shape),
            key=lambda shape: (int(shape.top), int(shape.left)),
        )
        for shape in shapes:
            if getattr(shape, "has_table", False):
                lines.extend(_table_rows(shape.table.rows))
            elif getattr(shape, "has_text_frame", False):
                value = shape.text.strip()
                if value:
                    lines.append(value)
        if slide.has_notes_slide:
            frame = slide.notes_slide.notes_text_frame
            notes = frame.text.strip() if frame is not None else ""
            if notes:
                lines.extend(("### Speaker notes", notes))
        slides.add("\n\n".join(lines))
    text = slides.text().strip()
    if not text:
        raise ExtractionError("PPTX contains no readable text")
    return ExtractionResult(
        text=text,
        media_type=schema.MEDIA_PPTX,
        extractor="python-pptx",
        pages=len(presentation.slides),
    )


def _image(
    data: bytes,
    media_type: str,
    *,
    ocr_model: str,
    vision_ocr: Callable[[bytes, str], str] | None,
) -> ExtractionResult:
    from PIL import Image

    try:
        with warnings.catch_warnings():
            warnings.simplefilter("error", Image.DecompressionBombWarning)
            with Image.open(io.BytesIO(data)) as image:
                width, height = image.size
    except (OSError, ValueError, Image.DecompressionBombWarning, Image.DecompressionBombError) as error:
        raise ArtifactRejected("image is corrupt or exceeds the pixel safety limit") from error
    if width <= 0 or height <= 0 or width * height > MAX_IMAGE_PIXELS:
        raise ArtifactRejected("image exceeds the pixel safety limit")

    try:
        text = _vision_handler(ocr_model, vision_ocr)(data, media_type).strip()
    except ArtifactRejected:
        raise
    except Exception as error:
        raise ExtractionError(f"image OCR failed ({error.__class__.__name__})") from error
    if not text:
        raise ExtractionError("image OCR found no readable text")
    return ExtractionResult(
        text=text,
        media_type=media_type,
        extractor=ocr_model.removeprefix("openrouter:"),
        ocr_pages=(1,),
        decisions=("image: OCR",),
    )


def _vision_handler(
    ocr_model: str,
    vision_ocr: Callable[[bytes, str], str] | None,
) -> Callable[[bytes, str], str]:
    if vision_ocr is not None:
        return vision_ocr
    from stigmergy.capture.vision import transcribe_image

    return lambda image, media: transcribe_image(image, media, model_name=ocr_model)


def _validate_ocr_model(ocr_model: str) -> None:
    if ocr_model != OCR_MODEL:
        raise ExtractionError(f"OCR model must be {OCR_MODEL}")


def _guard_pdf_ocr_geometry(width_points: float, height_points: float) -> None:
    if (
        not math.isfinite(width_points)
        or not math.isfinite(height_points)
        or width_points <= 0
        or height_points <= 0
    ):
        raise ArtifactRejected("PDF page has invalid geometry")
    width = math.ceil(width_points * OCR_DPI / 72)
    height = math.ceil(height_points * OCR_DPI / 72)
    if width * height > MAX_IMAGE_PIXELS:
        raise ArtifactRejected("PDF page exceeds the OCR pixel safety limit")


def _slack(data: bytes) -> ExtractionResult:
    from stigmergy.slack.snapshot import render_snapshot, validate_snapshot

    snapshot = validate_snapshot(data)
    return ExtractionResult(
        text=render_snapshot(snapshot),
        media_type=schema.MEDIA_SLACK,
        extractor="slack-snapshot",
        pages=1,
    )


def _clean_blocks(text: str) -> str:
    lines = [" ".join(line.split()) for line in text.splitlines()]
    compact = []
    for line in lines:
        if line or (compact and compact[-1]):
            compact.append(line)
    return "\n".join(compact).strip()


class _ReadableBuilder:
    def __init__(self) -> None:
        self.parts: list[str] = []
        self.bytes = 0

    def add(self, value: str) -> None:
        separator = "\n\n" if self.parts else ""
        self.bytes += len(separator.encode("utf-8")) + len(value.encode("utf-8"))
        if self.bytes > MAX_CAPTURE_EXTRACTED_BYTES:
            raise ExtractionError("extraction exceeds the readable byte limit")
        self.parts.append(value)

    def text(self) -> str:
        return "\n\n".join(self.parts)


def _guard_readable(value: str) -> None:
    if len(value.encode("utf-8")) > MAX_CAPTURE_EXTRACTED_BYTES:
        raise ExtractionError("extraction exceeds the readable byte limit")


def _apply_resource_limits(timeout_s: int) -> None:
    if not sys.platform.startswith("linux"):
        return
    import resource

    cpu_seconds = max(1, int(timeout_s))
    try:
        resource.setrlimit(
            resource.RLIMIT_AS,
            (MAX_CHILD_ADDRESS_SPACE_BYTES, MAX_CHILD_ADDRESS_SPACE_BYTES),
        )
        resource.setrlimit(resource.RLIMIT_CPU, (cpu_seconds, cpu_seconds + 1))
    except (OSError, ValueError) as error:
        raise ExtractionError("extractor resource limits could not be applied") from error


def extraction_clock() -> str:
    return dt.datetime.now(dt.UTC).isoformat()
